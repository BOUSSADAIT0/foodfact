"""Generate FoodFact presentation PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = r"c:\Users\dell\OneDrive\Bureau\propre_projet\foodfact\PRESENTATION_FoodFact_OpenData.pptx"

GREEN = RGBColor(0x22, 0xC5, 0x5E)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_title(slide, text, subtitle=None):
    title = slide.shapes.title
    title.text = text
    for p in title.text_frame.paragraphs:
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = DARK
    if subtitle and len(slide.placeholders) > 1:
        sub = slide.placeholders[1]
        sub.text = subtitle
        for p in sub.text_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY


def add_bullets(slide, items, left=0.6, top=1.6, width=8.8, height=5.0, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)


def add_section_header(slide, text):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()
    set_title(slide, text)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(
        slide,
        "FoodFact Recherche",
        "Application web & Open Data alimentaire\nProjet NutriRecherche — OpenFoodFacts",
    )


def slide_plan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "Plan de la présentation")
    add_bullets(
        slide,
        [
            "1. Contexte et objectifs du projet",
            "2. Qu'est-ce que l'Open Data ?",
            "3. OpenFoodFacts : notre source de données ouvertes",
            "4. Avons-nous utilisé de l'Open Data ? (Oui !)",
            "5. Architecture de l'application",
            "6. Technologies et outils utilisés",
            "7. Fonctionnalités de l'application",
            "8. Flux de données (Open Data → App)",
            "9. Données consommées et API",
            "10. Intérêt, limites et perspectives",
        ],
        font_size=17,
    )


def slide_contexte(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "Contexte et objectifs")
    add_bullets(
        slide,
        [
            "Problème : difficile de comparer rapidement les produits alimentaires (nutrition, scores, allergènes…)",
            "Objectif : créer une application web accessible pour rechercher et explorer des produits alimentaires",
            "Approche : réutiliser des données ouvertes existantes plutôt que créer une base propriétaire",
            "Public visé : consommateurs, étudiants, curieux de l'alimentation responsable",
            "Nom du projet : FoodFact Recherche (NutriRecherche)",
        ],
    )


def slide_open_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "Qu'est-ce que l'Open Data ?")
    add_bullets(
        slide,
        [
            "Données ouvertes = données librement accessibles, réutilisables et redistribuables",
            "Caractéristiques : gratuité, format ouvert (JSON, CSV…), licence permissive",
            "Intérêt : transparence, innovation, création d'applications sans recréer les données",
            "Exemples : OpenStreetMap (cartes), data.gouv.fr (administration), OpenFoodFacts (alimentation)",
            "Dans notre projet : nous consommons de l'Open Data, nous ne produisons pas de jeu de données",
        ],
    )


def slide_off(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "OpenFoodFacts — Base de données ouverte")
    add_bullets(
        slide,
        [
            "OpenFoodFacts (OFF) : base collaborative mondiale de produits alimentaires",
            "Site : https://world.openfoodfacts.org",
            "Licence : Open Database License (ODbL) — données ouvertes",
            "Contenu : nom, marque, ingrédients, allergènes, Nutri-Score, Éco-Score, NOVA, images…",
            "Alimentée par des bénévoles et des scans de codes-barres (crowdsourcing)",
            "Plus de 3 millions de produits référencés dans le monde",
        ],
    )


def slide_oui_open_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "Avons-nous fait de l'Open Data ?")
    add_bullets(
        slide,
        [
            "✅ OUI — Nous utilisons activement l'Open Data via OpenFoodFacts",
            "Type d'usage : CONSOMMATEUR d'Open Data (API REST publique)",
            "Nous ne publions pas notre propre jeu de données, mais une application qui valorise OFF",
            "API utilisée directement par le backend Scala :",
            "   • Recherche : world.openfoodfacts.org/cgi/search.pl (JSON)",
            "   • Détail produit : world.openfoodfacts.org/api/v0/product/{code}.json",
            "Lien retour vers OFF : bouton « Voir sur OpenFoodFacts » sur chaque fiche produit",
        ],
        font_size=16,
    )


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "Architecture de l'application")
    add_bullets(
        slide,
        [
            "Architecture client-serveur en 3 couches :",
            "",
            "🖥️ Frontend (Next.js) — port 3000",
            "   Interface utilisateur : recherche, filtres, graphiques, fiche produit",
            "",
            "⚙️ Backend (Scala / http4s) — port 8080",
            "   API REST intermédiaire : filtrage, tri, cache, pagination",
            "",
            "🌍 OpenFoodFacts (Open Data externe)",
            "   Source de vérité des données produits",
        ],
        font_size=16,
    )


def slide_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "Technologies utilisées")
    add_bullets(
        slide,
        [
            "Frontend :",
            "   • Next.js 16 + React 19 + TypeScript",
            "   • Tailwind CSS (interface responsive)",
            "   • Recharts (statistiques), Lucide (icônes)",
            "",
            "Backend :",
            "   • Scala 3.3 + sbt",
            "   • http4s (serveur HTTP), Circe (JSON), Cats Effect (async)",
            "",
            "DevOps / outils :",
            "   • Git, npm, Docker Compose, MUnit (tests)",
        ],
        font_size=15,
    )


def slide_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "Fonctionnalités principales")
    add_bullets(
        slide,
        [
            "🔍 Recherche textuelle avec debounce et URL partageable",
            "🏷️ Filtres : marque, pays, Nutri-Score (A–E), NOVA (1–4), calories/sucres/graisses",
            "📊 3 vues : grille, tableau triable, statistiques (graphiques)",
            "📄 Fiche produit : ingrédients, allergènes, scores, lien OpenFoodFacts",
            "🔄 Alternatives plus saines (même catégorie, tri Nutri-Score)",
            "📱 Pagination « Charger plus », mode clair/sombre",
        ],
        font_size=16,
    )


def slide_flux(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "Flux de données Open Data")
    add_bullets(
        slide,
        [
            "1. L'utilisateur saisit une recherche ou un filtre dans le navigateur",
            "2. Le frontend (Next.js) appelle notre API : GET /api/search?q=...",
            "3. Le backend Scala interroge OpenFoodFacts (JSON)",
            "4. OFF renvoie les produits bruts (Open Data)",
            "5. Notre backend filtre, trie, met en cache et renvoie une réponse enrichie",
            "6. Le frontend affiche cartes, tableaux et graphiques",
            "",
            "→ L'Open Data circule de OFF vers l'utilisateur via notre couche applicative",
        ],
        font_size=15,
    )


def slide_donnees(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "Données Open Data consommées")
    add_bullets(
        slide,
        [
            "Identité : code-barres, nom, marque, catégories, quantité, pays",
            "Nutrition (pour 100g) : énergie, sucres, graisses, sel, protéines, fibres",
            "Scores : Nutri-Score (A–E), Éco-Score, groupe NOVA (1–4)",
            "Santé : ingrédients, allergènes, labels (bio, sans gluten…)",
            "Médias : images produit (URL)",
            "",
            "Champs demandés explicitement à l'API OFF (optimisation des requêtes)",
        ],
        font_size=16,
    )


def slide_api(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "Notre API REST (couche intermédiaire)")
    add_bullets(
        slide,
        [
            "GET /health — état du serveur",
            "GET /api/search — recherche avec filtres et pagination",
            "GET /api/product/{code} — détail + alternatives",
            "",
            "Valeur ajoutée par rapport à OFF brut :",
            "   • Filtres combinés (pays, marque, Nutri-Score, NOVA, nutrition)",
            "   • Normalisation des pays (France, Tunisie, USA…)",
            "   • Cache mémoire, rate limiting, gestion d'erreurs",
            "   • Réponse structurée pour le frontend",
        ],
        font_size=15,
    )


def slide_interet(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "Intérêt, limites et perspectives")
    add_bullets(
        slide,
        [
            "Intérêt Open Data :",
            "   • Pas de base à construire manuellement",
            "   • Données riches, mises à jour par la communauté",
            "   • Projet pédagogique sur la réutilisation de données ouvertes",
            "",
            "Limites :",
            "   • Qualité variable selon les produits (données incomplètes)",
            "   • Dépendance à l'API externe (latence, disponibilité)",
            "",
            "Perspectives : export CSV, comparaison produits, scan code-barres mobile",
        ],
        font_size=15,
    )


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_section_header(slide, "Conclusion")
    add_bullets(
        slide,
        [
            "FoodFact Recherche est une application web qui VALORISE l'Open Data alimentaire",
            "Source Open Data : OpenFoodFacts (licence ouverte, API publique)",
            "Stack : Next.js (frontend) + Scala/http4s (backend)",
            "L'Open Data permet de se concentrer sur l'expérience utilisateur et l'analyse",
            "",
            "Merci pour votre attention !",
            "Questions ?",
        ],
        font_size=18,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_plan(prs)
    slide_contexte(prs)
    slide_open_data(prs)
    slide_off(prs)
    slide_oui_open_data(prs)
    slide_architecture(prs)
    slide_tech(prs)
    slide_features(prs)
    slide_flux(prs)
    slide_donnees(prs)
    slide_api(prs)
    slide_interet(prs)
    slide_conclusion(prs)

    prs.save(OUTPUT)
    print(f"Presentation saved: {OUTPUT}")


if __name__ == "__main__":
    main()
